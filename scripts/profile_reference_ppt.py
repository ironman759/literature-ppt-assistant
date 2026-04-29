#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from collections import Counter
from pathlib import Path

from pptx import Presentation


def rgb_tuple(rgb) -> list[int] | None:
    if rgb is None:
        return None
    return [rgb[0], rgb[1], rgb[2]]


def iter_text_shapes(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if text:
                yield shape, text


def profile(path: Path) -> dict:
    prs = Presentation(path)
    slide_profiles = []
    font_counter: Counter[str] = Counter()
    color_counter: Counter[tuple[int, int, int]] = Counter()
    notes_count = 0

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        shape_count = len(slide.shapes)
        picture_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:
                picture_count += 1
        for shape, text in iter_text_shapes(slide):
            texts.append(text)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        font_counter[run.font.name] += 1
                    if run.font.color and run.font.color.rgb:
                        color = rgb_tuple(run.font.color.rgb)
                        if color:
                            color_counter[tuple(color)] += 1
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes:
            notes_count += 1
        slide_profiles.append(
            {
                "slide_no": idx,
                "text_blocks": len(texts),
                "char_count": sum(len(t) for t in texts),
                "picture_count": picture_count,
                "shape_count": shape_count,
                "title_guess": texts[0][:120] if texts else "",
                "has_notes": bool(notes),
            }
        )

    avg_chars = round(sum(s["char_count"] for s in slide_profiles) / max(len(slide_profiles), 1), 1)
    return {
        "file": str(path),
        "slide_count": len(slide_profiles),
        "average_chars_per_slide": avg_chars,
        "slides_with_notes": notes_count,
        "dominant_fonts": [name for name, _ in font_counter.most_common(5)],
        "dominant_text_colors": [list(color) for color, _ in color_counter.most_common(5)],
        "slide_profiles": slide_profiles,
        "style_guidance": {
            "pace": "match reference slide count and text density when scientifically appropriate",
            "notes": "match note presence and oral-delivery depth when the reference has notes",
            "visual_density": "use picture_count and text_blocks as guardrails, not hard constraints",
        },
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Profile a reference PPTX for literature presentation pacing and style.")
    parser.add_argument("--pptx", required=True, type=Path, help="Reference PPTX")
    parser.add_argument("--output", required=True, type=Path, help="Output JSON profile")
    args = parser.parse_args()

    payload = profile(args.pptx)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    print(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
