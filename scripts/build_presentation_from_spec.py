#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_THEME = {
    "palette": {
        "background": [248, 245, 239],
        "primary": [26, 77, 90],
        "secondary": [181, 137, 78],
        "text": [33, 37, 41],
        "muted": [92, 100, 107],
        "surface": [255, 255, 255],
        "line": [226, 236, 234],
        "band": [239, 232, 219],
        "accent": [180, 45, 45],
    },
    "typography": {
        "title_font": "Georgia",
        "body_font": "Georgia",
        "label_font": "Georgia",
    },
    "shape": {
        "radius": "rounded",
        "border_width": 1.0,
    },
}


def merge_tokens(base: dict, override: dict | None) -> dict:
    merged = json.loads(json.dumps(base))
    if not override:
        return merged
    for key, value in override.items():
        if isinstance(value, dict) and isinstance(merged.get(key), dict):
            merged[key].update(value)
        else:
            merged[key] = value
    return merged


def theme_for(spec: dict, slide_spec: dict | None = None) -> dict:
    theme = merge_tokens(DEFAULT_THEME, spec.get("design_tokens"))
    if slide_spec:
        theme = merge_tokens(theme, slide_spec.get("design_tokens"))
    return theme


def color(theme: dict, name: str) -> RGBColor:
    rgb = theme["palette"][name]
    return RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def is_accent_rgb(rgb: RGBColor | None, theme: dict) -> bool:
    if rgb is None:
        return False
    accent = color(theme, "accent")
    return tuple(rgb) == tuple(accent)


def font_name(theme: dict, role: str) -> str:
    return theme["typography"].get(role, DEFAULT_THEME["typography"][role])


def set_background(slide, theme: dict) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(theme, "background")


def add_background_accents(slide, theme: dict) -> None:
    bottom_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.08), Inches(13.33), Inches(0.42)
    )
    bottom_band.fill.solid()
    bottom_band.fill.fore_color.rgb = color(theme, "band")
    bottom_band.line.fill.background()


def highlight_token(token: str) -> bool:
    if not token.strip():
        return False
    patterns = [
        r"\b\d+(?:\.\d+)?%?\b",
        r"\b\d+(?:–|-)\d+\b",
        r"\b(?:MiFish(?:-[A-Za-z]+)?|eDNA|PCR|MiSeq|BLAST|FASTQ|ND5|COI|12S|mtDNA)\b",
    ]
    return any(re.search(pattern, token) for pattern in patterns)


def add_rich_text(paragraph, text: str, font_size: int, theme: dict, inverse: bool = False) -> None:
    segments = re.split(r"(\*\*.*?\*\*)", text)
    for segment in segments:
        if not segment:
            continue
        is_marked = segment.startswith("**") and segment.endswith("**")
        clean = segment[2:-2] if is_marked else segment
        parts = re.split(r"(\s+)", clean)
        for part in parts:
            if part == "":
                continue
            run = paragraph.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.name = font_name(theme, "body_font")
            if is_marked or highlight_token(part):
                run.font.bold = True
                run.font.color.rgb = color(theme, "accent")
            else:
                run.font.color.rgb = color(theme, "surface") if inverse else color(theme, "text")


def add_top_bar(
    slide,
    section: str,
    title: str,
    subtitle: str | None = None,
    *,
    assertion_title: str | None = None,
    slide_purpose: str | None = None,
    theme: dict,
) -> None:
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(0.34), Inches(0.72), Inches(0.44)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = color(theme, "primary")
    chip.line.fill.background()

    sec_box = slide.shapes.add_textbox(Inches(0.54), Inches(0.43), Inches(0.45), Inches(0.18))
    p = sec_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = color(theme, "surface")

    title_box = slide.shapes.add_textbox(Inches(1.28), Inches(0.2), Inches(9.25), Inches(0.72))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.line_spacing = 0.95
    run = p.add_run()
    run.text = assertion_title or title
    run.font.size = Pt(21)
    run.font.bold = True
    run.font.color.rgb = color(theme, "primary")
    run.font.name = font_name(theme, "title_font")

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.82), Inches(10.4), Inches(0.2))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(10.5)
        run.font.color.rgb = color(theme, "muted")
        run.font.name = font_name(theme, "body_font")

    if slide_purpose:
        purpose_box = slide.shapes.add_textbox(Inches(10.55), Inches(0.27), Inches(2.0), Inches(0.25))
        p = purpose_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = slide_purpose
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = color(theme, "secondary")
        run.font.name = font_name(theme, "label_font")

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), Inches(1.12), Inches(12.15), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color(theme, "line")
    line.line.fill.background()


def add_bullets(slide, bullets: list[str], has_image: bool = False, *, theme: dict) -> None:
    panel_width = Inches(7.45 if has_image else 11.95)
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.48), Inches(1.18), panel_width, Inches(5.35)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = color(theme, "surface")
    panel.line.color.rgb = color(theme, "line")
    panel.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

    box = slide.shapes.add_textbox(Inches(0.78), Inches(1.47), Inches(6.85 if has_image else 11.2), Inches(4.7))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        bullet_run = p.add_run()
        bullet_run.text = "• "
        bullet_run.font.size = Pt(18 if has_image else 20)
        bullet_run.font.color.rgb = color(theme, "secondary")
        bullet_run.font.name = font_name(theme, "body_font")
        add_rich_text(p, bullet, 18 if has_image else 20, theme)
        p.space_after = Pt(12)
        p.line_spacing = 1.15


def add_image_panel(slide, image_path: str | None, caption: str | None = None, *, theme: dict) -> None:
    if not image_path:
        return

    path = Path(image_path)
    if not path.exists():
        return

    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.22), Inches(1.18), Inches(4.6), Inches(5.35)
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = color(theme, "surface")
    frame.line.color.rgb = color(theme, "secondary")
    frame.line.width = Pt(1.2)

    label_box = slide.shapes.add_textbox(Inches(8.48), Inches(1.34), Inches(1.7), Inches(0.2))
    p = label_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Paper Figure"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = color(theme, "secondary")

    slide.shapes.add_picture(str(path), Inches(8.48), Inches(1.66), width=Inches(4.08), height=Inches(3.52))

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(8.48), Inches(5.3), Inches(4.02), Inches(0.78))
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = caption
        run.font.size = Pt(10)
        run.font.color.rgb = color(theme, "muted")


def add_takeaway(slide, text: str | None, *, theme: dict, label: str = "Takeaway") -> None:
    if not text:
        return

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.48), Inches(6.04), Inches(12.28), Inches(0.6)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color(theme, "band")
    card.line.fill.background()

    label_box = slide.shapes.add_textbox(Inches(0.72), Inches(6.18), Inches(1.2), Inches(0.18))
    p = label_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = color(theme, "primary")
    run.font.name = font_name(theme, "label_font")

    msg_box = slide.shapes.add_textbox(Inches(1.95), Inches(6.16), Inches(10.4), Inches(0.2))
    p = msg_box.text_frame.paragraphs[0]
    add_rich_text(p, text, 12, theme)


def add_section_intro(slide, slide_spec: dict, spec: dict) -> None:
    theme = theme_for(spec, slide_spec)
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.55), Inches(11.5), Inches(4.55)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color(theme, "surface")
    card.line.color.rgb = color(theme, "line")
    card.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(2.0), Inches(0.22), Inches(2.8)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color(theme, "primary")
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.62), Inches(2.05), Inches(7.8), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_spec.get("assertion_title") or slide_spec.get("title", "")
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = color(theme, "primary")
    run.font.name = font_name(theme, "title_font")

    subtitle = slide_spec.get("subtitle")
    if subtitle:
        box = slide.shapes.add_textbox(Inches(1.66), Inches(3.26), Inches(7.1), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(16)
        run.font.color.rgb = color(theme, "text")
        run.font.name = font_name(theme, "body_font")

    bullets = slide_spec.get("bullets", [])
    if bullets:
        box = slide.shapes.add_textbox(Inches(1.68), Inches(4.1), Inches(6.9), Inches(1.45))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for idx, bullet in enumerate(bullets[:3]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = ""
            dot = p.add_run()
            dot.text = "• "
            dot.font.size = Pt(18)
            dot.font.color.rgb = color(theme, "secondary")
            dot.font.name = font_name(theme, "body_font")
            add_rich_text(p, bullet, 17, theme)
            p.space_after = Pt(10)

    takeaway = slide_spec.get("takeaway") or slide_spec.get("slide_purpose")
    if takeaway:
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.02), Inches(2.12), Inches(2.4), Inches(1.05)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = color(theme, "band")
        pill.line.fill.background()
        box = slide.shapes.add_textbox(Inches(9.26), Inches(2.36), Inches(1.95), Inches(0.56))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = takeaway
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = color(theme, "primary")
        run.font.name = font_name(theme, "label_font")


def add_workflow_layout(slide, slide_spec: dict, spec: dict) -> None:
    theme = theme_for(spec, slide_spec)
    steps = slide_spec.get("workflow_steps") or slide_spec.get("bullets", [])
    steps = steps[:4]
    if not steps:
        render_standard_layout(slide, slide_spec, spec)
        return

    total = len(steps)
    left = 0.72
    width = 2.8 if total >= 4 else 3.25
    gap = 0.18
    top = 2.05
    height = 2.45

    for idx, step in enumerate(steps):
        x = Inches(left + idx * (width + gap))
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color(theme, "surface")
        card.line.color.rgb = color(theme, "line")
        card.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + idx * (width + gap) + 0.2), Inches(top + 0.18), Inches(0.5), Inches(0.5)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = color(theme, "primary")
        badge.line.fill.background()

        num_box = slide.shapes.add_textbox(
            Inches(left + idx * (width + gap) + 0.31), Inches(top + 0.28), Inches(0.28), Inches(0.16)
        )
        p = num_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(idx + 1)
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = color(theme, "surface")
        run.font.name = font_name(theme, "label_font")

        box = slide.shapes.add_textbox(
            Inches(left + idx * (width + gap) + 0.28), Inches(top + 0.82), Inches(width - 0.45), Inches(1.35)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        add_rich_text(p, step, 15, theme)

        if idx < total - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(left + idx * (width + gap) + width - 0.02),
                Inches(top + 0.95),
                Inches(0.24),
                Inches(0.42),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color(theme, "secondary")
            arrow.line.fill.background()

    add_takeaway(slide, slide_spec.get("takeaway"), theme=theme, label="Workflow")


def add_result_summary_layout(slide, slide_spec: dict, spec: dict) -> None:
    theme = theme_for(spec, slide_spec)
    highlights = slide_spec.get("result_highlights", [])
    if not highlights:
        bullets = slide_spec.get("bullets", [])
        for bullet in bullets[:3]:
            text = bullet.replace("**", "")
            number_match = re.search(r"(\d+(?:\.\d+)?%?)", text)
            label = text
            value = number_match.group(1) if number_match else "Key"
            if number_match:
                label = text.replace(value, "").strip(" ,:-")
            highlights.append({"value": value, "label": label or text})
    highlights = highlights[:3]
    if not highlights:
        render_standard_layout(slide, slide_spec, spec)
        return

    card_y = 2.0
    width = 3.45
    gap = 0.32
    for idx, item in enumerate(highlights):
        x = 0.82 + idx * (width + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(card_y), Inches(width), Inches(2.15)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color(theme, "surface")
        card.line.color.rgb = color(theme, "line")
        card.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

        value_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(card_y + 0.28), Inches(width - 0.44), Inches(0.72))
        p = value_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(item.get("value", "Key"))
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = color(theme, "primary")
        run.font.name = font_name(theme, "title_font")

        label_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(card_y + 1.08), Inches(width - 0.44), Inches(0.72))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        add_rich_text(p, str(item.get("label", "")), 14, theme)

    summary_bullets = slide_spec.get("bullets", [])[:3]
    if summary_bullets:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(4.55), Inches(10.85), Inches(1.25)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = color(theme, "surface")
        panel.line.color.rgb = color(theme, "line")
        panel.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

        box = slide.shapes.add_textbox(Inches(1.08), Inches(4.82), Inches(10.3), Inches(0.78))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for idx, bullet in enumerate(summary_bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = ""
            dot = p.add_run()
            dot.text = "• "
            dot.font.size = Pt(14)
            dot.font.color.rgb = color(theme, "secondary")
            dot.font.name = font_name(theme, "body_font")
            add_rich_text(p, bullet, 13, theme)
            p.space_after = Pt(8)
    add_takeaway(slide, slide_spec.get("takeaway"), theme=theme, label="Result")


def add_comparison_layout(slide, slide_spec: dict, spec: dict) -> None:
    theme = theme_for(spec, slide_spec)
    left = slide_spec.get("left_column", {})
    right = slide_spec.get("right_column", {})

    positions = [
        (0.62, 1.5, 5.95, 4.85, left, color(theme, "primary")),
        (6.77, 1.5, 5.95, 4.85, right, color(theme, "secondary")),
    ]

    for x, y, w, h, column, bar_color in positions:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color(theme, "surface")
        card.line.color.rgb = color(theme, "line")
        card.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.26)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        header_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.4), Inches(w - 0.44), Inches(0.38))
        p = header_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(column.get("title", "Comparison"))
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = color(theme, "text")
        run.font.name = font_name(theme, "title_font")

        subtitle = column.get("subtitle")
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.78), Inches(w - 0.44), Inches(0.28))
            p = sub_box.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(subtitle)
            run.font.size = Pt(10.5)
            run.font.color.rgb = color(theme, "muted")
            run.font.name = font_name(theme, "body_font")

        bullets = column.get("bullets", [])
        if bullets:
            box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 1.15), Inches(w - 0.5), Inches(2.2))
            tf = box.text_frame
            tf.word_wrap = True
            tf.clear()
            for idx, bullet in enumerate(bullets[:4]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = ""
                dot = p.add_run()
                dot.text = "• "
                dot.font.size = Pt(14)
                dot.font.color.rgb = bar_color
                dot.font.name = font_name(theme, "body_font")
                add_rich_text(p, bullet, 13, theme)
                p.space_after = Pt(8)

        image_path = column.get("image_path")
        if image_path and Path(image_path).exists():
            frame_y = y + 3.35
            slide.shapes.add_picture(str(image_path), Inches(x + 0.28), Inches(frame_y), width=Inches(w - 0.56), height=Inches(1.35))

    add_takeaway(slide, slide_spec.get("takeaway"), theme=theme, label="Compare")


def add_limitation_implication_layout(slide, slide_spec: dict, spec: dict) -> None:
    theme = theme_for(spec, slide_spec)
    blocks = [
        ("Limitations", slide_spec.get("limitations", []), 0.72, color(theme, "accent")),
        ("Implications", slide_spec.get("implications", []), 6.82, color(theme, "primary")),
    ]
    for title, items, x, accent in blocks:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(5.8), Inches(4.95)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color(theme, "surface")
        card.line.color.rgb = color(theme, "line")
        card.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

        header = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.24), Inches(1.8), Inches(2.0), Inches(0.46)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = accent
        header.line.fill.background()
        box = slide.shapes.add_textbox(Inches(x + 0.45), Inches(1.93), Inches(1.6), Inches(0.18))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = color(theme, "surface")
        run.font.name = font_name(theme, "label_font")

        body = slide.shapes.add_textbox(Inches(x + 0.28), Inches(2.48), Inches(5.2), Inches(3.55))
        tf = body.text_frame
        tf.word_wrap = True
        tf.clear()
        for idx, item in enumerate(items[:5]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = ""
            dot = p.add_run()
            dot.text = "• "
            dot.font.size = Pt(15)
            dot.font.color.rgb = accent
            dot.font.name = font_name(theme, "body_font")
            add_rich_text(p, item, 14, theme)
            p.space_after = Pt(10)

    add_takeaway(slide, slide_spec.get("takeaway"), theme=theme, label="Interpret")


def add_figure_focus_layout(slide, slide_spec: dict, spec: dict) -> None:
    theme = theme_for(spec, slide_spec)
    image_path = slide_spec.get("image_path")
    if not image_path or not Path(image_path).exists():
        render_standard_layout(slide, slide_spec, spec)
        return

    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.66), Inches(1.38), Inches(7.25), Inches(4.95)
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = color(theme, "surface")
    frame.line.color.rgb = color(theme, "secondary")
    frame.line.width = Pt(1.2)
    slide.shapes.add_picture(str(image_path), Inches(0.9), Inches(1.64), width=Inches(6.78), height=Inches(4.1))

    caption = slide_spec.get("image_caption")
    if caption:
        box = slide.shapes.add_textbox(Inches(0.92), Inches(5.87), Inches(6.72), Inches(0.32))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = caption
        run.font.size = Pt(9.5)
        run.font.color.rgb = color(theme, "muted")
        run.font.name = font_name(theme, "body_font")

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.18), Inches(1.38), Inches(4.55), Inches(4.95)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = color(theme, "surface")
    panel.line.color.rgb = color(theme, "line")
    panel.line.width = Pt(float(theme["shape"].get("border_width", 1.0)))

    callout = slide_spec.get("figure_callout")
    if callout:
        callout_box = slide.shapes.add_textbox(Inches(8.48), Inches(1.72), Inches(3.95), Inches(0.95))
        tf = callout_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(callout)
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = color(theme, "primary")
        run.font.name = font_name(theme, "title_font")

    bullets = slide_spec.get("bullets", [])
    if bullets:
        box = slide.shapes.add_textbox(Inches(8.48), Inches(2.82), Inches(3.9), Inches(2.6))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for idx, bullet in enumerate(bullets[:4]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = ""
            dot = p.add_run()
            dot.text = "• "
            dot.font.size = Pt(14)
            dot.font.color.rgb = color(theme, "secondary")
            dot.font.name = font_name(theme, "body_font")
            add_rich_text(p, bullet, 13, theme)
            p.space_after = Pt(8)

    add_takeaway(slide, slide_spec.get("takeaway"), theme=theme, label="Figure")


def add_footer(slide, index: int, citation: str | None, *, theme: dict) -> None:
    num_box = slide.shapes.add_textbox(Inches(12.55), Inches(7.15), Inches(0.28), Inches(0.16))
    p = num_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(index)
    run.font.size = Pt(8.5)
    run.font.color.rgb = color(theme, "muted")

    if citation:
        cite_box = slide.shapes.add_textbox(Inches(0.54), Inches(7.13), Inches(8.6), Inches(0.16))
        p = cite_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = citation
        run.font.size = Pt(8.5)
        run.font.color.rgb = color(theme, "muted")


def add_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


def add_title_slide(prs: Presentation, spec: dict) -> None:
    theme = theme_for(spec)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, theme)
    add_background_accents(slide, theme)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(0.84), Inches(0.24), Inches(4.85)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color(theme, "primary")
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.22), Inches(1.2), Inches(9.8), Inches(1.95))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec["deck_title"]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = color(theme, "primary")
    run.font.name = font_name(theme, "title_font")

    subtitle = spec.get("subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.25), Inches(3.42), Inches(8.5), Inches(0.65))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(18)
        run.font.color.rgb = color(theme, "text")
        run.font.name = font_name(theme, "body_font")

    info_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.18), Inches(4.96), Inches(5.85), Inches(1.05)
    )
    info_card.fill.solid()
    info_card.fill.fore_color.rgb = color(theme, "surface")
    info_card.line.color.rgb = color(theme, "line")
    info_card.line.width = Pt(1.0)

    citation = spec.get("citation")
    presenter = spec.get("presenter", "Prepared by Codex")
    footer_box = slide.shapes.add_textbox(Inches(1.42), Inches(5.19), Inches(5.2), Inches(0.6))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = citation or ""
    run.font.size = Pt(15)
    run.font.color.rgb = color(theme, "text")
    run.font.name = font_name(theme, "body_font")
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = presenter
    run2.font.size = Pt(12)
    run2.font.color.rgb = color(theme, "muted")
    run2.font.name = font_name(theme, "body_font")

    add_notes(slide, spec.get("title_notes"))


def add_sampling_design_layout(slide, slide_spec: dict) -> None:
    theme = theme_for({}, slide_spec)
    subtitle = slide_spec.get("subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.42), Inches(1.18), Inches(5.2), Inches(0.34))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(18)
        run.font.color.rgb = color(theme, "muted")
        run.font.name = font_name(theme, "body_font")

    left_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.62), Inches(6.35), Inches(3.68)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = color(theme, "surface")
    left_card.line.color.rgb = color(theme, "line")
    left_card.line.width = Pt(1.0)

    right_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.62), Inches(5.95), Inches(3.68)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = color(theme, "surface")
    right_card.line.color.rgb = color(theme, "line")
    right_card.line.width = Pt(1.0)

    # Accent bars
    left_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(1.62), Inches(0.04), Inches(3.68)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = color(theme, "primary")
    left_bar.line.fill.background()

    right_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.0), Inches(1.62), Inches(0.04), Inches(3.68)
    )
    right_bar.fill.solid()
    right_bar.fill.fore_color.rgb = color(theme, "secondary")
    right_bar.line.fill.background()

    headings = slide_spec.get("panel_headings", {})
    left_head = slide.shapes.add_textbox(Inches(0.72), Inches(1.78), Inches(2.8), Inches(0.3))
    p = left_head.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = headings.get("left", "Aquarium Tanks")
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = color(theme, "primary")
    run.font.name = font_name(theme, "title_font")

    right_head = slide.shapes.add_textbox(Inches(7.22), Inches(1.78), Inches(3.2), Inches(0.3))
    p = right_head.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = headings.get("right", "Natural Seawater")
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = color(theme, "primary")
    run.font.name = font_name(theme, "title_font")

    left_img = Path(slide_spec.get("left_image_path", ""))
    if left_img.exists():
        slide.shapes.add_picture(str(left_img), Inches(0.74), Inches(2.15), width=Inches(5.75), height=Inches(2.45))
    right_img = Path(slide_spec.get("right_image_path", ""))
    if right_img.exists():
        slide.shapes.add_picture(str(right_img), Inches(7.23), Inches(2.15), width=Inches(5.42), height=Inches(2.45))

    left_cap = slide.shapes.add_textbox(Inches(0.78), Inches(4.72), Inches(5.6), Inches(0.3))
    p = left_cap.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = slide_spec.get("left_caption", "Fig. 1a-d Aquarium tanks used for validation")
    run.font.size = Pt(10)
    run.font.color.rgb = color(theme, "muted")
    run.font.name = font_name(theme, "body_font")

    right_cap = slide.shapes.add_textbox(Inches(7.28), Inches(4.72), Inches(5.2), Inches(0.3))
    p = right_cap.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = slide_spec.get("right_caption", "Fig. 1e-f Coral reef site and field map")
    run.font.size = Pt(10)
    run.font.color.rgb = color(theme, "muted")
    run.font.name = font_name(theme, "body_font")

    # Bottom advantage bars
    left_adv = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(5.45), Inches(6.35), Inches(0.9)
    )
    left_adv.fill.solid()
    left_adv.fill.fore_color.rgb = color(theme, "primary")
    left_adv.line.fill.background()

    right_adv = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(5.45), Inches(5.95), Inches(0.9)
    )
    right_adv.fill.solid()
    right_adv.fill.fore_color.rgb = color(theme, "secondary")
    right_adv.line.fill.background()

    left_adv_text = slide.shapes.add_textbox(Inches(0.78), Inches(5.63), Inches(5.6), Inches(0.5))
    tf = left_adv_text.text_frame
    p = tf.paragraphs[0]
    add_rich_text(
        p,
        slide_spec.get("left_advantage", "**Known species composition** allows **accuracy validation** and a clear expected answer."),
        11,
        theme,
        inverse=True,
    )
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.color.rgb = color(theme, "accent") if is_accent_rgb(run.font.color.rgb, theme) else color(theme, "surface")
            run.font.name = font_name(theme, "body_font")

    right_adv_text = slide.shapes.add_textbox(Inches(7.28), Inches(5.63), Inches(5.2), Inches(0.5))
    tf = right_adv_text.text_frame
    p = tf.paragraphs[0]
    add_rich_text(
        p,
        slide_spec.get("right_advantage", "Unknown community composition tests **practical applicability** in real-world scenarios."),
        11,
        theme,
        inverse=True,
    )
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.color.rgb = color(theme, "accent") if is_accent_rgb(run.font.color.rgb, theme) else color(theme, "surface")
            run.font.name = font_name(theme, "body_font")

    add_takeaway(slide, slide_spec.get("takeaway"), theme=theme)


def render_standard_layout(slide, slide_spec: dict, spec: dict) -> None:
    theme = theme_for(spec, slide_spec)
    has_image = bool(slide_spec.get("image_path"))
    add_bullets(slide, slide_spec.get("bullets", []), has_image=has_image, theme=theme)
    add_image_panel(slide, slide_spec.get("image_path"), slide_spec.get("image_caption"), theme=theme)
    add_takeaway(slide, slide_spec.get("takeaway"), theme=theme)


LAYOUT_RENDERERS = {
    "comparison": add_comparison_layout,
    "figure_focus": add_figure_focus_layout,
    "limitation_implication": add_limitation_implication_layout,
    "section_intro": add_section_intro,
    "workflow": add_workflow_layout,
    "result_summary": add_result_summary_layout,
    "standard": render_standard_layout,
    "sampling_design": add_sampling_design_layout,
}


def render_slide_layout(slide, slide_spec: dict, spec: dict) -> None:
    layout = slide_spec.get("layout", "standard")
    renderer = LAYOUT_RENDERERS.get(layout, render_standard_layout)
    if renderer is add_sampling_design_layout:
        merged_slide_spec = dict(slide_spec)
        merged_slide_spec["design_tokens"] = merge_tokens(spec.get("design_tokens", {}), slide_spec.get("design_tokens"))
        add_sampling_design_layout(slide, merged_slide_spec)
        return
    renderer(slide, slide_spec, spec)


def build_presentation(spec: dict, output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, spec)

    for index, slide_spec in enumerate(spec["slides"], start=2):
        theme = theme_for(spec, slide_spec)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, theme)
        add_background_accents(slide, theme)
        if slide_spec.get("layout", "standard") != "section_intro":
            add_top_bar(
                slide,
                slide_spec["section"],
                slide_spec["title"],
                slide_spec.get("subtitle"),
                assertion_title=slide_spec.get("assertion_title"),
                slide_purpose=slide_spec.get("slide_purpose"),
                theme=theme,
            )
        render_slide_layout(slide, slide_spec, spec)
        add_footer(slide, index - 1, spec.get("citation"), theme=theme)
        add_notes(slide, slide_spec.get("notes"))

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a literature presentation PPTX from a JSON spec.")
    parser.add_argument("--spec", required=True, type=Path, help="Path to JSON spec file")
    parser.add_argument("--output", required=True, type=Path, help="Output PPTX path")
    args = parser.parse_args()

    with args.spec.open("r", encoding="utf-8") as f:
        spec = json.load(f)

    build_presentation(spec, args.output)
    print(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
